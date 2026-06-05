"""
文档处理与向量存储模块
负责文件解析、文本切分、向量化存储和检索
"""
import hashlib
import json
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings, AzureOpenAIEmbeddings
from langchain_community.vectorstores import FAISS

import config


def _get_embeddings():
    """根据配置返回 Embedding 模型"""
    if config.LLM_PROVIDER == "azure":
        return AzureOpenAIEmbeddings(
            azure_endpoint=config.AZURE_OPENAI_ENDPOINT,
            api_key=config.AZURE_OPENAI_API_KEY,
            azure_deployment=config.EMBEDDING_MODEL,
            api_version=config.AZURE_OPENAI_API_VERSION,
        )
    return OpenAIEmbeddings(
        api_key=config.OPENAI_API_KEY,
        base_url=config.OPENAI_BASE_URL,
        model=config.EMBEDDING_MODEL,
    )


def parse_file(file_path: str) -> str:
    """解析上传的文件，返回纯文本内容"""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".txt" or suffix == ".md":
        return path.read_text(encoding="utf-8")

    if suffix == ".pdf":
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages)

    if suffix == ".docx":
        from docx import Document
        doc = Document(str(path))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n\n".join(texts)

    raise ValueError(f"不支持的文件类型: {suffix}")


class DocumentStore:
    """文档向量存储，支持增量添加和相似度检索"""

    def __init__(self):
        self.embeddings = _get_embeddings()
        self.vector_store: FAISS | None = None
        self.doc_metadata: dict[str, dict] = {}  # file_hash -> metadata
        self._meta_path = config.DATA_DIR / "doc_metadata.json"
        self._load_metadata()
        self._try_load_vector_store()

    # ----------------------------------------------------------
    # 持久化
    # ----------------------------------------------------------

    def _load_metadata(self):
        if self._meta_path.exists():
            self.doc_metadata = json.loads(self._meta_path.read_text("utf-8"))

    def _save_metadata(self):
        self._meta_path.write_text(
            json.dumps(self.doc_metadata, ensure_ascii=False, indent=2), "utf-8"
        )

    def _try_load_vector_store(self):
        index_path = config.VECTOR_STORE_DIR / "index.faiss"
        if index_path.exists():
            self.vector_store = FAISS.load_local(
                str(config.VECTOR_STORE_DIR),
                self.embeddings,
                allow_dangerous_deserialization=True,
            )

    def _save_vector_store(self):
        if self.vector_store:
            self.vector_store.save_local(str(config.VECTOR_STORE_DIR))

    # ----------------------------------------------------------
    # 文档操作
    # ----------------------------------------------------------

    def add_document(self, file_path: str, file_name: str) -> dict:
        """解析文件并加入向量存储，返回文档摘要信息"""
        text = parse_file(file_path)
        file_hash = hashlib.sha256(text.encode()).hexdigest()[:16]

        if file_hash in self.doc_metadata:
            return self.doc_metadata[file_hash]

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.CHUNK_SIZE,
            chunk_overlap=config.CHUNK_OVERLAP,
            separators=["\n\n", "\n", "。", "；", " "],
        )
        chunks = splitter.split_text(text)

        metadatas = [{"source": file_name, "file_hash": file_hash, "chunk_idx": i}
                     for i in range(len(chunks))]

        if self.vector_store is None:
            self.vector_store = FAISS.from_texts(chunks, self.embeddings, metadatas=metadatas)
        else:
            self.vector_store.add_texts(chunks, metadatas=metadatas)

        doc_info = {
            "file_name": file_name,
            "file_hash": file_hash,
            "total_chunks": len(chunks),
            "total_chars": len(text),
            "full_text": text[:5000],  # 保留前 5000 字作为摘要依据
        }
        self.doc_metadata[file_hash] = doc_info

        self._save_vector_store()
        self._save_metadata()
        return doc_info

    def search(self, query: str, top_k: int = 5) -> list[dict]:
        """根据查询进行相似度检索"""
        if not self.vector_store:
            return []
        results = self.vector_store.similarity_search_with_score(query, k=top_k)
        return [
            {"content": doc.page_content, "source": doc.metadata.get("source", ""), "score": float(score)}
            for doc, score in results
        ]

    def get_all_documents(self) -> list[dict]:
        """返回所有已入库文档的元数据"""
        return list(self.doc_metadata.values())

    def get_combined_context(self, max_chars: int = 8000) -> str:
        """拼接所有文档摘要文本，用于给其他 Agent 提供完整知识背景"""
        texts = []
        total = 0
        for info in self.doc_metadata.values():
            snippet = info.get("full_text", "")
            if total + len(snippet) > max_chars:
                snippet = snippet[: max_chars - total]
            texts.append(f"【{info['file_name']}】\n{snippet}")
            total += len(snippet)
            if total >= max_chars:
                break
        return "\n\n---\n\n".join(texts)

    def clear(self):
        """清空所有存储"""
        self.vector_store = None
        self.doc_metadata = {}
        self._save_metadata()
        for p in config.VECTOR_STORE_DIR.iterdir():
            p.unlink()
